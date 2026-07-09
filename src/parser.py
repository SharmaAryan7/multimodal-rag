import fitz  # PyMuPDF
import pdfplumber
import os
import csv
import hashlib
import tempfile
import atexit
import shutil
from PIL import Image
from io import BytesIO
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import markdown
from src.config import (
    DATA_DIR,
    CHUNK_SIZE, CHUNK_OVERLAP,
    MIN_TEXT_LENGTH, MIN_IMAGE_SIZE
)


SUPPORTED_EXTENSIONS = (
    ".pdf", ".docx",
    ".txt", ".md",
    ".html", ".htm",
    ".pptx", ".xlsx", ".csv",
    ".png", ".jpg", ".jpeg", ".bmp", ".gif"
)


# Temporary folder (auto-deleted on exit)
_TEMP_DIR = tempfile.mkdtemp(prefix="mmrag_imgs_")

# Track seen image hashes for deduplication
_seen_image_hashes = set()


def _cleanup_temp_dir():
    try:
        shutil.rmtree(_TEMP_DIR, ignore_errors=True)
    except Exception:
        pass


atexit.register(_cleanup_temp_dir)


def reset_dedup_cache():
    """Reset deduplication cache. Call before each new ingestion run."""
    global _seen_image_hashes
    _seen_image_hashes = set()


def extract_text_chunks(text, page_num, source_file):
    chunks = []
    if not text or len(text.strip()) < MIN_TEXT_LENGTH:
        return chunks

    text = text.strip()
    start = 0
    chunk_index = 0

    while start < len(text):
        end = start + CHUNK_SIZE
        chunk_text = text[start:end]

        if len(chunk_text.strip()) >= MIN_TEXT_LENGTH:
            chunks.append({
                "text": chunk_text.strip(),
                "content_type": "text",
                "page": page_num,
                "source_file": source_file,
                "chunk_index": chunk_index
            })
            chunk_index += 1

        start += CHUNK_SIZE - CHUNK_OVERLAP

    return chunks


def classify_image(image_path):
    """Classify image as 'flowchart' or 'image' using strict heuristics."""
    try:
        img = Image.open(image_path).convert("RGB")
        width, height = img.size

        pixels = list(img.getdata())
        total_pixels = len(pixels)
        sample_size = min(10000, total_pixels)
        step = max(1, total_pixels // sample_size)
        sampled = pixels[::step]

        white_count = sum(1 for r, g, b in sampled if r > 240 and g > 240 and b > 240)
        white_ratio = white_count / len(sampled)

        quantized = set()
        for r, g, b in sampled:
            quantized.add((r // 30, g // 30, b // 30))
        num_colors = len(quantized)

        gradient_count = 0
        sharp_count = 0
        img_arr = list(img.getdata())
        check_step = max(2, width // 80)

        for y in range(10, height - 10, max(2, height // 60)):
            for x in range(10, width - 10, check_step):
                idx = y * width + x
                idx_right = idx + check_step

                if idx_right < total_pixels:
                    r1, g1, b1 = img_arr[idx]
                    r2, g2, b2 = img_arr[idx_right]

                    diff = abs(r1 - r2) + abs(g1 - g2) + abs(b1 - b2)

                    if 5 < diff < 40:
                        gradient_count += 1
                    elif diff > 150:
                        sharp_count += 1

        total_checks = gradient_count + sharp_count + 1
        gradient_ratio = gradient_count / total_checks

        is_flowchart = (
            white_ratio > 0.50
            and num_colors < 100
            and gradient_ratio < 0.4
        )

        return "flowchart" if is_flowchart else "image"

    except Exception as e:
        print(f"   Classification error: {e}")
        return "image"


def save_temp_image(image_bytes, label="img"):
    """Save image to temp folder with deduplication via hash."""
    min_w = MIN_IMAGE_SIZE[0]
    min_h = MIN_IMAGE_SIZE[1]

    try:
        # Compute MD5 hash to detect duplicates
        img_hash = hashlib.md5(image_bytes).hexdigest()

        if img_hash in _seen_image_hashes:
            return None, None, 0, 0  # Skip silently — duplicate

        pil_image = Image.open(BytesIO(image_bytes))
        width, height = pil_image.size

        if width < min_w or height < min_h:
            return None, None, width, height

        _seen_image_hashes.add(img_hash)

        tmp = tempfile.NamedTemporaryFile(
            suffix=".jpeg",
            prefix=f"{label}_",
            dir=_TEMP_DIR,
            delete=False
        )
        tmp_path = tmp.name
        tmp.close()

        pil_image = pil_image.convert("RGB")
        pil_image.save(tmp_path, "JPEG", quality=90)

        content_type = classify_image(tmp_path)
        return tmp_path, content_type, width, height

    except Exception as e:
        print(f"      Error saving image: {e}")
        return None, None, 0, 0


# ==================== PDF ====================

def extract_images_from_pdf(pdf_path, source_file):
    print("\n  STAGE: IMAGE EXTRACTION (PDF)")
    print("-" * 60)

    images = []
    skipped_duplicates = 0
    doc = fitz.open(pdf_path)

    for page_num in range(len(doc)):
        page = doc[page_num]
        image_list = page.get_images(full=True)

        if image_list:
            print(f"  Page {page_num + 1}: Found {len(image_list)} image(s)")

        for img_index, img in enumerate(image_list):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]

                label = f"p{page_num + 1}_i{img_index + 1}"
                img_path, content_type, w, h = save_temp_image(image_bytes, label)

                if not img_path:
                    if w == 0 and h == 0:
                        skipped_duplicates += 1
                    else:
                        print(f"      Image {img_index + 1}: SKIPPED (too small: {w}x{h})")
                    continue

                print(f"      Image {img_index + 1}: {w}x{h} -> {content_type.upper()}")

                images.append({
                    "path": img_path,
                    "content_type": content_type,
                    "page": page_num + 1,
                    "source_file": source_file
                })

            except Exception as e:
                print(f"      Error: {e}")

    doc.close()

    if skipped_duplicates > 0:
        print(f"  Deduplication: skipped {skipped_duplicates} duplicate image(s)")

    return images


def extract_tables_from_pdf(pdf_path, source_file):
    print("\n  STAGE: TABLE EXTRACTION (PDF)")
    print("-" * 60)

    table_chunks = []
    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages):
            tables = page.extract_tables()
            for t_idx, table in enumerate(tables):
                if not table or len(table) < 2:
                    continue

                print(f"  Page {page_num + 1}, Table {t_idx + 1}: {len(table)} rows x {len(table[0])} cols")

                md_lines = []
                headers = [str(h).replace("\n", " ").strip() if h else "" for h in table[0]]
                md_lines.append("| " + " | ".join(headers) + " |")
                md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

                for row in table[1:]:
                    clean_row = [str(c).replace("\n", " ").strip() if c else "" for c in row]
                    md_lines.append("| " + " | ".join(clean_row) + " |")

                markdown_table = "\n".join(md_lines)

                if len(markdown_table) >= MIN_TEXT_LENGTH:
                    table_chunks.append({
                        "text": f"[TABLE on page {page_num + 1}]\n{markdown_table}",
                        "content_type": "table",
                        "page": page_num + 1,
                        "source_file": source_file,
                        "chunk_index": 0
                    })

    return table_chunks


def parse_pdf(pdf_path, source_file):
    text_chunks = []
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        text = doc[page_num].get_text()
        text_chunks.extend(extract_text_chunks(text, page_num + 1, source_file))
    doc.close()

    images = extract_images_from_pdf(pdf_path, source_file)
    text_chunks.extend(extract_tables_from_pdf(pdf_path, source_file))
    return text_chunks, images


# ==================== DOCX ====================

def parse_docx(docx_path, source_file):
    print("\n  STAGE: DOCX PARSING")
    print("-" * 60)

    doc = Document(docx_path)
    full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    text_chunks = extract_text_chunks(full_text, 1, source_file)

    for t_idx, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])

        if len(rows) < 2:
            continue

        md_lines = ["| " + " | ".join(rows[0]) + " |"]
        md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
        for row in rows[1:]:
            md_lines.append("| " + " | ".join(row) + " |")

        markdown_table = "\n".join(md_lines)
        if len(markdown_table) >= MIN_TEXT_LENGTH:
            text_chunks.append({
                "text": f"[TABLE]\n{markdown_table}",
                "content_type": "table",
                "page": 1,
                "source_file": source_file,
                "chunk_index": t_idx
            })

    images = []
    img_count = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                image_bytes = rel.target_part.blob
                img_count += 1
                label = f"docx_i{img_count}"
                img_path, content_type, w, h = save_temp_image(image_bytes, label)

                if img_path:
                    print(f"  Image {img_count}: {w}x{h} -> {content_type.upper()}")
                    images.append({
                        "path": img_path,
                        "content_type": content_type,
                        "page": 1,
                        "source_file": source_file
                    })
            except Exception as e:
                print(f"  Error: {e}")

    return text_chunks, images


# ==================== TXT / MD ====================

def parse_txt(txt_path, source_file):
    print("\n  STAGE: TXT/MD PARSING")
    print("-" * 60)

    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()

    if source_file.endswith(".md"):
        html = markdown.markdown(content)
        content = BeautifulSoup(html, "html.parser").get_text()

    print(f"  Loaded {len(content)} characters")
    text_chunks = extract_text_chunks(content, 1, source_file)
    return text_chunks, []


# ==================== HTML ====================

def parse_html(html_path, source_file):
    print("\n  STAGE: HTML PARSING")
    print("-" * 60)

    with open(html_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    for tag in soup(["script", "style"]):
        tag.decompose()

    content = soup.get_text(separator="\n")
    content = "\n".join([line.strip() for line in content.splitlines() if line.strip()])

    print(f"  Extracted {len(content)} characters")
    text_chunks = extract_text_chunks(content, 1, source_file)

    for t_idx, table in enumerate(soup.find_all("table")):
        rows = []
        for row in table.find_all("tr"):
            cells = [c.get_text(strip=True) for c in row.find_all(["td", "th"])]
            if cells:
                rows.append(cells)

        if len(rows) < 2:
            continue

        md_lines = ["| " + " | ".join(rows[0]) + " |"]
        md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
        for row in rows[1:]:
            md_lines.append("| " + " | ".join(row) + " |")

        markdown_table = "\n".join(md_lines)
        if len(markdown_table) >= MIN_TEXT_LENGTH:
            text_chunks.append({
                "text": f"[TABLE]\n{markdown_table}",
                "content_type": "table",
                "page": 1,
                "source_file": source_file,
                "chunk_index": t_idx
            })

    return text_chunks, []


# ==================== PPTX ====================

def parse_pptx(pptx_path, source_file):
    print("\n  STAGE: PPTX PARSING")
    print("-" * 60)

    prs = Presentation(pptx_path)
    text_chunks = []
    images = []
    img_count = 0

    for slide_num, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text.append(paragraph.text.strip())

            if shape.shape_type == 13:
                try:
                    image = shape.image
                    img_count += 1
                    label = f"slide{slide_num + 1}_i{img_count}"
                    img_path, content_type, w, h = save_temp_image(image.blob, label)
                    if img_path:
                        print(f"  Slide {slide_num + 1}, Image: {w}x{h} -> {content_type.upper()}")
                        images.append({
                            "path": img_path,
                            "content_type": content_type,
                            "page": slide_num + 1,
                            "source_file": source_file
                        })
                except Exception as e:
                    print(f"  Error: {e}")

        slide_content = "\n".join(slide_text)
        if slide_content:
            text_chunks.extend(extract_text_chunks(slide_content, slide_num + 1, source_file))

    print(f"  Processed {len(prs.slides)} slide(s)")
    return text_chunks, images


# ==================== XLSX ====================

def parse_xlsx(xlsx_path, source_file):
    print("\n  STAGE: XLSX PARSING")
    print("-" * 60)

    wb = load_workbook(xlsx_path, data_only=True)
    text_chunks = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cleaned = [str(c) if c is not None else "" for c in row]
            if any(cell.strip() for cell in cleaned):
                rows.append(cleaned)

        if len(rows) < 2:
            continue

        print(f"  Sheet '{sheet_name}': {len(rows)} rows x {len(rows[0])} cols")

        md_lines = ["| " + " | ".join(rows[0]) + " |"]
        md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
        for row in rows[1:]:
            md_lines.append("| " + " | ".join(row) + " |")

        markdown_table = "\n".join(md_lines)
        if len(markdown_table) >= MIN_TEXT_LENGTH:
            text_chunks.append({
                "text": f"[SHEET: {sheet_name}]\n{markdown_table}",
                "content_type": "table",
                "page": sheet_idx + 1,
                "source_file": source_file,
                "chunk_index": 0
            })

    return text_chunks, []


# ==================== CSV ====================

def parse_csv(csv_path, source_file):
    print("\n  STAGE: CSV PARSING")
    print("-" * 60)

    rows = []
    with open(csv_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append([c.strip() for c in row])

    if len(rows) < 2:
        return [], []

    print(f"  CSV: {len(rows)} rows x {len(rows[0])} cols")

    md_lines = ["| " + " | ".join(rows[0]) + " |"]
    md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    for row in rows[1:]:
        md_lines.append("| " + " | ".join(row) + " |")

    markdown_table = "\n".join(md_lines)
    text_chunks = []
    if len(markdown_table) >= MIN_TEXT_LENGTH:
        text_chunks.append({
            "text": f"[CSV TABLE]\n{markdown_table}",
            "content_type": "table",
            "page": 1,
            "source_file": source_file,
            "chunk_index": 0
        })

    return text_chunks, []


# ==================== IMAGE FILE ====================

def parse_image(image_path, source_file):
    print("\n  STAGE: IMAGE FILE PARSING")
    print("-" * 60)

    try:
        with open(image_path, "rb") as f:
            image_bytes = f.read()

        label = "imgfile"
        img_path, content_type, w, h = save_temp_image(image_bytes, label)

        if not img_path:
            print(f"  SKIPPED (too small or duplicate)")
            return [], []

        print(f"  Image: {w}x{h} -> {content_type.upper()}")

        return [], [{
            "path": img_path,
            "content_type": content_type,
            "page": 1,
            "source_file": source_file
        }]
    except Exception as e:
        print(f"  Error: {e}")
        return [], []


# ==================== DISPATCHER ====================

def parse_file(filepath, source_file):
    ext = os.path.splitext(source_file)[1].lower()

    print(f"\n{'=' * 60}")
    print(f"  PARSING: {source_file} ({ext})")
    print(f"{'=' * 60}")

    if ext == ".pdf":
        return parse_pdf(filepath, source_file)
    elif ext == ".docx":
        return parse_docx(filepath, source_file)
    elif ext in (".txt", ".md"):
        return parse_txt(filepath, source_file)
    elif ext in (".html", ".htm"):
        return parse_html(filepath, source_file)
    elif ext == ".pptx":
        return parse_pptx(filepath, source_file)
    elif ext == ".xlsx":
        return parse_xlsx(filepath, source_file)
    elif ext == ".csv":
        return parse_csv(filepath, source_file)
    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".gif"):
        return parse_image(filepath, source_file)
    else:
        print(f"  Unsupported format: {ext}")
        return [], []


def parse_all_documents():
    all_text_chunks = []
    all_images = []

    # Reset dedup cache for fresh ingestion
    reset_dedup_cache()

    if not os.path.exists(DATA_DIR):
        print("No data directory found.")
        return all_text_chunks, all_images

    files = [f for f in os.listdir(DATA_DIR) if f.lower().endswith(SUPPORTED_EXTENSIONS)]

    if not files:
        print(f"No supported files found. Supported: {SUPPORTED_EXTENSIONS}")
        return all_text_chunks, all_images

    print(f"\nFound {len(files)} file(s) to process")

    for filename in files:
        filepath = os.path.join(DATA_DIR, filename)
        try:
            text_chunks, images = parse_file(filepath, filename)
            all_text_chunks.extend(text_chunks)
            all_images.extend(images)
        except Exception as e:
            print(f"   Error parsing {filename}: {e}")

    return all_text_chunks, all_images